# -*- coding: utf-8 -*-
"""공시→데이터→기반Tool→종합Tool(오케스트레이터) 5컬럼 매핑을 PPTX + 미리보기 SVG로 생성.
Anthropic/Claude 화이트 테마. orchestration(tool of tools)을 5번째 컬럼으로 통합."""
import html
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# ---------- 노드 ----------
L0 = [("ty1","정기보고서"),("ty2","주요사항보고서"),("ty3","거래소 수시공시"),
      ("ty4","대량보유·지분 보고"),("ty5","공시 메타")]
L1 = [("periodic","사업·반기·분기보고서"),("d_ts","자기주식 취득·처분·신탁 결정"),
      ("d_issue","유상증자·CB·BW·감자 결정"),("d_merge","합병·분할·주식교환 결정"),
      ("d_deal","타법인주식 양수·양도 결정"),("d_riskB","회생·부도·영업정지·해산"),
      ("d_notice","주주총회 소집공고"),("d_result","주주총회 결과(의결권행사)"),
      ("d_tsresult","자기주식 결과보고서·소각"),("d_divdec","현금·현물배당 결정"),
      ("d_order","단일판매·공급계약 체결"),("d_gov","기업지배구조보고서"),
      ("d_valueup","기업가치제고계획(밸류업)"),("d_riskI","중대재해·횡령배임·생산중단"),
      ("d_lawsuit","경영권분쟁소송 제기"),("d_5pct","5% 대량보유상황보고서"),
      ("d_ownchg","최대주주 소유주식 변동신고"),("d_proxy","의결권 대리행사 권유(위임장)"),
      ("d_tender","공개매수 신고서"),("d_list","전체 공시 목록"),("d_recept","rcept_no 메타")]
L2 = [("da_fs","재무제표(BS·IS·CF)"),("da_audit","감사의견"),("da_divtable","배당에 관한 사항"),
      ("da_majorholder","최대주주·특수관계인"),("da_shares","주식총수·자기주식 현황"),
      ("da_tscond","취득·처분·신탁 조건"),("da_issuecond","발행조건·잠재희석률"),
      ("da_mergeratio","합병비율·외부평가"),("da_dealterm","상대방·금액·자산대비%"),
      ("da_riskB","사유·영향(도산·정지)"),("da_agenda","안건·이사후보·보수한도"),
      ("da_vote","안건별 가결/부결·찬반율"),("da_tsactual","실제 취득·처분·소각 수량"),
      ("da_dps","DPS·배당총액·기준일"),("da_contract","계약금액·매출대비%"),
      ("da_gov15","15개 지표 준수율"),("da_commit","commitment·이행"),
      ("da_riskI","사상자·혐의금액"),("da_lawsuit","소송 단계·당사자"),
      ("da_coholder","공동보유자 분해"),("da_5dyn","5% 동학(매집·목적)"),
      ("da_ownchg","최대주주 지분 변동"),("da_proxyrec","위임장 권유 내역"),
      ("da_tendercond","공개매수 조건"),("da_index","회사 식별·공시 인덱스"),
      ("da_viewer","공시일·소스·뷰어 URL")]
# 기반 tool (15) — proxy_advise/proxy_contest 는 5번째 컬럼으로 분리
L3 = [("t_fin","financial_metrics"),("t_div","dividend"),("t_own","ownership_structure"),
      ("t_tre","treasury_share"),("t_dil","dilutive_issuance"),("t_res","corporate_restructuring"),
      ("t_dea","corporate_deals"),("t_rsk","risk_events"),("t_not","shareholder_meeting_notice"),
      ("t_rslt","shareholder_meeting_results"),("t_ord","order_contracts"),
      ("t_gov","corp_gov_report"),("t_val","value_up"),("t_cmp","company"),("t_evd","evidence")]
# 종합 tool (오케스트레이터) + 법령 — 커스텀 배치
L4 = [
    ("law","상법 법령 레이어\n(wiki/rules/laws)"),
    ("t_adv","proxy_advise_before_meeting\n★ TOOL OF TOOLS ★\n내부: director_evaluation\n→ 안건별 의결권 권고"),
    ("t_pcn","proxy_contest\n→ 경영권 분쟁 신호\n(+ 소송·위임장·공개매수·5% 공시 직접조회)"),
]

# 일반 흐름 (공시→데이터→기반tool)
E_NORMAL = [
    ("ty1","periodic"),("ty2","d_ts"),("ty2","d_issue"),("ty2","d_merge"),("ty2","d_deal"),("ty2","d_riskB"),
    ("ty3","d_notice"),("ty3","d_result"),("ty3","d_tsresult"),("ty3","d_divdec"),("ty3","d_order"),
    ("ty3","d_gov"),("ty3","d_valueup"),("ty3","d_riskI"),("ty3","d_lawsuit"),
    ("ty4","d_5pct"),("ty4","d_ownchg"),("ty4","d_proxy"),("ty4","d_tender"),("ty5","d_list"),("ty5","d_recept"),
    ("periodic","da_fs"),("periodic","da_audit"),("periodic","da_divtable"),("periodic","da_majorholder"),("periodic","da_shares"),
    ("d_ts","da_tscond"),("d_issue","da_issuecond"),("d_merge","da_mergeratio"),("d_deal","da_dealterm"),("d_riskB","da_riskB"),
    ("d_notice","da_agenda"),("d_result","da_vote"),("d_tsresult","da_tsactual"),("d_divdec","da_dps"),("d_order","da_contract"),
    ("d_gov","da_gov15"),("d_valueup","da_commit"),("d_riskI","da_riskI"),("d_lawsuit","da_lawsuit"),
    ("d_5pct","da_coholder"),("d_5pct","da_5dyn"),("d_ownchg","da_ownchg"),("d_proxy","da_proxyrec"),
    ("d_tender","da_tendercond"),("d_list","da_index"),("d_recept","da_viewer"),
    ("da_fs","t_fin"),("da_audit","t_fin"),("da_divtable","t_div"),("da_dps","t_div"),
    ("da_majorholder","t_own"),("da_shares","t_own"),("da_coholder","t_own"),("da_ownchg","t_own"),
    ("da_tscond","t_tre"),("da_tsactual","t_tre"),("da_issuecond","t_dil"),("da_mergeratio","t_res"),
    ("da_dealterm","t_dea"),("da_riskB","t_rsk"),("da_riskI","t_rsk"),("da_agenda","t_not"),
    ("da_vote","t_rslt"),("da_contract","t_ord"),("da_gov15","t_gov"),("da_commit","t_val"),
    ("da_index","t_cmp"),("da_viewer","t_evd"),
]
# 오케스트레이션 (기반tool → 종합tool) — tool of tools
E_ORCH = [("t_own","t_pcn"),("t_rslt","t_pcn"),
          ("t_fin","t_adv"),("t_div","t_adv"),("t_own","t_adv"),("t_tre","t_adv"),
          ("t_ord","t_adv"),("t_gov","t_adv"),("t_not","t_adv")]
# proxy_contest 직접 공시 조회 (데이터 → 종합tool)
E_DIRECT = [("da_5dyn","t_pcn"),("da_lawsuit","t_pcn"),("da_proxyrec","t_pcn"),("da_tendercond","t_pcn")]
# 법령 → proxy_advise
E_LAW = [("law","t_adv")]

LAYERS = [L0, L1, L2, L3]
TITLES = ["① 공시 타입","② 공시 이름","③ 추출 데이터","④ 기반 Tool","⑤ 종합 Tool (오케스트레이터)"]
FILLS  = ["F3DCD2","DEE6EC","E3E8D8","F1E3CA"]
BORDERS= ["CC785C","7C93A8","8A9A6B","C09A5B"]
TITLE_COLS = ["BC5A38","5E7388","6E7E50","A87F3C","BC5A38"]
ORCH_FILL = "EFC9B8"; ORCH_BORDER = "CC785C"
LAW_FILL  = "F1E3CA"; LAW_BORDER = "C09A5B"
ORCH_LINE = "CC785C"; DIRECT_LINE = "9A8C7A"
BG = "FAF9F5"; TEXT = "2B2A28"; HEAD_TEXT = "1F1E1D"

# ---------- 지오메트리 ----------
SW, SH = 19.0, 10.0
COL_X = [0.30, 2.20, 6.30, 10.25, 14.35]
COL_W = [1.75, 3.85, 3.55, 3.60, 4.45]
CONTENT_TOP, CONTENT_BOT = 1.35, 0.30
USABLE = SH - CONTENT_TOP - CONTENT_BOT
BOX_H = 0.275

geo = {}
for li, layer in enumerate(LAYERS):
    n = len(layer); x = COL_X[li]; w = COL_W[li]
    for i,(nid,_) in enumerate(layer):
        cy = CONTENT_TOP + (i+0.5)*USABLE/n
        geo[nid] = (x, cy-BOX_H/2, w, BOX_H, li)
# L4 커스텀 (x, y_top, w, h, layer)
x4, w4 = COL_X[4], COL_W[4]
geo["law"]   = (x4, 1.75-0.27, w4, 0.55, 4)
geo["t_adv"] = (x4, 4.30-0.55, w4, 1.10, 4)
geo["t_pcn"] = (x4, 7.20-0.45, w4, 0.95, 4)

ALL_NODES = {nid: lbl for layer in LAYERS for nid,lbl in layer}
ALL_NODES.update({nid: lbl for nid,lbl in L4})

def node_style(nid, li):
    if nid in ("t_adv","t_pcn"): return ORCH_FILL, ORCH_BORDER
    if nid == "law": return LAW_FILL, LAW_BORDER
    return FILLS[li], BORDERS[li]

# ===================== PPTX =====================
prs = Presentation(); prs.slide_width = Inches(SW); prs.slide_height = Inches(SH)
slide = prs.slides.add_slide(prs.slide_layouts[6])
bgf = slide.background.fill; bgf.solid(); bgf.fore_color.rgb = RGBColor.from_string(BG)

def add_text(x,y,w,h,txt,size,color,bold=False,align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf = tb.text_frame; tf.word_wrap=True; p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=txt; f=r.font; f.size=Pt(size); f.bold=bold
    f.color.rgb=RGBColor.from_string(color); f.name="Malgun Gothic"

add_text(0.3,0.16,18,0.6,"OpenProxy MCP — 공시 → 데이터 → Tool → 종합 Tool 매핑",26,HEAD_TEXT,True)
for li,t in enumerate(TITLES):
    cnt = len(L4) if li==4 else len(LAYERS[li])
    add_text(COL_X[li],0.78,COL_W[li],0.45,f"{t} ({cnt})",14,TITLE_COLS[li],True,PP_ALIGN.CENTER)

def emu(v): return Emu(int(v*914400))
def add_conn(a,b,color,wpt,dash=False,alpha=55000):
    xa,ya,wa,ha,_=geo[a]; xb,yb,wb,hb,_=geo[b]
    cn=slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,emu(xa+wa),emu(ya+ha/2),emu(xb),emu(yb+hb/2))
    ln=cn.line; ln.color.rgb=RGBColor.from_string(color); ln.width=Pt(wpt)
    el=ln._get_or_add_ln(); sf=el.find(qn('a:solidFill'))
    if sf is not None:
        srgb=sf.find(qn('a:srgbClr'))
        srgb.append(srgb.makeelement(qn('a:alpha'),{'val':str(alpha)}))
    if dash:
        el.append(el.makeelement(qn('a:prstDash'),{'val':'dash'}))

for a,b in E_NORMAL:  add_conn(a,b,BORDERS[geo[a][4]],0.75,False,52000)
for a,b in E_ORCH:    add_conn(a,b,ORCH_LINE,1.5,False,90000)
for a,b in E_DIRECT:  add_conn(a,b,DIRECT_LINE,0.9,True,55000)
for a,b in E_LAW:     add_conn(a,b,LAW_BORDER,1.4,False,90000)

def add_box(nid,li):
    label=ALL_NODES[nid]; x,y,w,h,_=geo[nid]; fill,border=node_style(nid,li)
    sh=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=RGBColor.from_string(fill)
    sh.line.color.rgb=RGBColor.from_string(border); sh.line.width=Pt(2.0 if nid in ("t_adv","t_pcn") else 1.0)
    sh.shadow.inherit=False
    tf=sh.text_frame; tf.word_wrap=True
    tf.margin_left=Emu(20000); tf.margin_right=Emu(20000); tf.margin_top=Emu(4000); tf.margin_bottom=Emu(4000)
    tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    lines=label.split("\n")
    for idx,line in enumerate(lines):
        p=tf.paragraphs[0] if idx==0 else tf.add_paragraph()
        p.alignment=PP_ALIGN.CENTER
        r=p.add_run(); r.text=line; f=r.font
        big = nid in ("t_adv","t_pcn")
        f.size=Pt(9.5 if (big and idx==0) else (8 if big else 8.5))
        f.bold=(li>=3 or nid in ("t_adv","t_pcn")) and idx==0
        f.color.rgb=RGBColor.from_string(TEXT); f.name="Malgun Gothic"

for li,layer in enumerate(LAYERS):
    for nid,_ in layer: add_box(nid,li)
for nid,_ in L4: add_box(nid,4)

prs.save("wiki/tools/diagrams/tool_disclosure_map.pptx"); print("PPTX saved")

# ===================== 미리보기 SVG =====================
PX=96
def esc(s): return html.escape(s,quote=True)
Wp,Hp=int(SW*PX),int(SH*PX)
svg=[f'<svg xmlns="http://www.w3.org/2000/svg" width="{Wp}" height="{Hp}" '
     f'font-family="Malgun Gothic, sans-serif"><rect width="{Wp}" height="{Hp}" fill="#{BG}"/>']
svg.append(f'<text x="{int(0.3*PX)}" y="{int(0.52*PX)}" fill="#{HEAD_TEXT}" font-size="26" font-weight="700">'
           f'OpenProxy MCP — 공시 → 데이터 → Tool → 종합 Tool 매핑</text>')
for li,t in enumerate(TITLES):
    cx=(COL_X[li]+COL_W[li]/2)*PX; cnt=len(L4) if li==4 else len(LAYERS[li])
    svg.append(f'<text x="{cx:.0f}" y="{int(1.02*PX)}" fill="#{TITLE_COLS[li]}" font-size="15" '
               f'font-weight="700" text-anchor="middle">{esc(t)} ({cnt})</text>')
def sline(a,b,color,w,op,dash=False):
    xa,ya,wa,ha,_=geo[a]; xb,yb,wb,hb,_=geo[b]
    x1=(xa+wa)*PX; y1=(ya+ha/2)*PX; x2=xb*PX; y2=(yb+hb/2)*PX; mx=(x1+x2)/2
    d=' stroke-dasharray="5 4"' if dash else ''
    svg.append(f'<path d="M{x1:.1f},{y1:.1f} C{mx:.1f},{y1:.1f} {mx:.1f},{y2:.1f} {x2:.1f},{y2:.1f}" '
               f'fill="none" stroke="#{color}" stroke-width="{w}" stroke-opacity="{op}"{d}/>')
for a,b in E_NORMAL: sline(a,b,BORDERS[geo[a][4]],1.0,0.45)
for a,b in E_DIRECT: sline(a,b,DIRECT_LINE,1.1,0.5,True)
for a,b in E_ORCH:   sline(a,b,ORCH_LINE,1.8,0.8)
for a,b in E_LAW:    sline(a,b,LAW_BORDER,1.6,0.85)
def srect(nid,li):
    label=ALL_NODES[nid]; x,y,w,h,_=geo[nid]; px,py,pw,ph=x*PX,y*PX,w*PX,h*PX
    fill,border=node_style(nid,li); big=nid in ("t_adv","t_pcn")
    svg.append(f'<rect x="{px:.1f}" y="{py:.1f}" width="{pw:.1f}" height="{ph:.1f}" rx="5" '
               f'fill="#{fill}" stroke="#{border}" stroke-width="{2 if big else 1.2}"/>')
    lines=label.split("\n"); n=len(lines); fs=10.5 if big else 11
    y0=py+ph/2-(n-1)*(fs+2)/2+4
    for i,ln in enumerate(lines):
        fw=700 if (li>=3 or big) and i==0 else 400
        svg.append(f'<text x="{px+pw/2:.1f}" y="{y0+i*(fs+2):.1f}" fill="#{TEXT}" font-size="{fs}" '
                   f'font-weight="{fw}" text-anchor="middle">{esc(ln)}</text>')
for li,layer in enumerate(LAYERS):
    for nid,_ in layer: srect(nid,li)
for nid,_ in L4: srect(nid,4)
svg.append('</svg>')
with open("wiki/tools/diagrams/ppt_preview.svg","w",encoding="utf-8") as f: f.write("\n".join(svg))
print(f"SVG preview written: {Wp}x{Hp}")
